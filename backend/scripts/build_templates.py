r"""Render the branded .pptx templates into ``backend/templates/`` from the
definitions in the repo-root ``templates.json``.

Python port of the retired Node ``scripts/build-templates.ts`` (pptxgenjs).
Every placeholder is created as a rectangle auto-shape whose **shape name**
equals its ``shape`` field in ``templates.json`` so Office.js can find it by
``shape.name`` (the same contract ``apply_slide_content`` fills against). The
generated ``.pptx`` files are committed artifacts.

Run with:  cd backend ; .\.venv\Scripts\python scripts\build_templates.py
Optional:  --out <dir>   write elsewhere (default: backend/templates/)
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# backend/scripts/build_templates.py -> repo root is two parents up.
_REPO_ROOT = Path(__file__).resolve().parent.parent.parent
_TEMPLATES_JSON = _REPO_ROOT / "templates.json"
_DEFAULT_OUT = Path(__file__).resolve().parent.parent / "templates"

# Widescreen canvas (matches pptxgenjs LAYOUT_WIDE: 13.333in x 7.5in).
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
CANVAS_W = 13.33  # coordinate width used by the layout math (as in the TS original)

# Brand palette (generic "navy/gold" house style).
NAVY = "1F3864"
NAVY_LIGHT = "2E5395"
GOLD = "C9A227"
GRAY_TEXT = "595959"
CARD_BG = "F2F4F8"
WHITE = "FFFFFF"
FONT = "Segoe UI"


def _set_shape_name(shape, name: str) -> None:
    """Set the OOXML cNvPr/@name so Office.js reads it back as ``shape.name``."""
    shape._element.nvSpPr.cNvPr.name = name


def _add_rect(slide, x, y, w, h, *, fill, line=None, line_w=None, rounded=False, radius=None):
    """A decorative rectangle (brand bar or KPI card). No text, no name."""
    shape_kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded and radius is not None:
        # pptxgenjs rectRadius is an absolute inch radius; the OOXML adjustment is a
        # fraction of the shorter side. Keep the same subtle corner.
        try:
            sp.adjustments[0] = max(0.0, min(0.5, radius / min(w, h)))
        except (IndexError, ValueError):
            pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is not None:
        sp.line.color.rgb = RGBColor.from_string(line)
        if line_w is not None:
            sp.line.width = Pt(line_w)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def _add_text(
    slide,
    name,
    text,
    x,
    y,
    w,
    h,
    *,
    size,
    color,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    shrink=False,
):
    """A named text placeholder: a transparent rectangle carrying text, exactly
    like pptxgenjs ``addText`` with an ``objectName``."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    # Transparent: no fill, no outline, no shadow — text only.
    sp.fill.background()
    sp.line.fill.background()
    sp.shadow.inherit = False
    _set_shape_name(sp, name)

    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if shrink else MSO_AUTO_SIZE.NONE
    # Zero the insets so text sits where pptxgenjs put it.
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = RGBColor.from_string(color)
    return sp


def _blank_slide(prs):
    """A truly empty slide (layout 6 = "Blank" in the default template)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)


def build_quarterly_review(prs: Presentation) -> None:
    # ---- Slide 1: Title ----------------------------------------------------
    title = _blank_slide(prs)
    _set_bg(title, NAVY)
    _add_rect(title, 0, 0, CANVAS_W, 0.25, fill=GOLD)          # top brand bar
    _add_rect(title, 0.9, 2.5, 2.2, 0.07, fill=GOLD)           # accent rule
    _add_text(title, "title_text", "Presentation Title",
              0.9, 2.7, 11.5, 1.4, size=44, bold=True, color=WHITE)
    _add_text(title, "subtitle_text", "Client / engagement subtitle",
              0.9, 4.1, 11.5, 0.8, size=22, color="D6DCE5")
    _add_text(title, "date_text", "Date",
              0.9, 6.5, 6, 0.5, size=14, color=GOLD)

    # ---- Slide 2: KPI Summary ----------------------------------------------
    kpi = _blank_slide(prs)
    _set_bg(kpi, WHITE)
    _add_rect(kpi, 0, 0, CANVAS_W, 1.1, fill=NAVY)             # header banner
    _add_rect(kpi, 0, 1.1, CANVAS_W, 0.06, fill=GOLD)          # banner underline
    _add_text(kpi, "heading_text", "Financial Highlights",
              0.6, 0, 12.1, 1.1, size=26, bold=True, color=WHITE, shrink=True)

    # Four KPI cards in a centred row.
    card_w, gap = 2.9, 0.25
    row_x = (CANVAS_W - (4 * card_w + 3 * gap)) / 2
    for i in range(1, 5):
        x = row_x + (i - 1) * (card_w + gap)
        _add_rect(kpi, x, 2.2, card_w, 2.6, fill=CARD_BG, line=NAVY_LIGHT, line_w=0.75,
                  rounded=True, radius=0.08)
        _add_text(kpi, f"kpi{i}_label", f"KPI {i} label",
                  x + 0.15, 2.5, card_w - 0.3, 0.5, size=14, bold=True, color=GRAY_TEXT,
                  align=PP_ALIGN.CENTER, shrink=True)
        _add_text(kpi, f"kpi{i}_value", "—",
                  x + 0.15, 3.2, card_w - 0.3, 1.1, size=28, bold=True, color=NAVY,
                  align=PP_ALIGN.CENTER, shrink=True)

    _add_text(kpi, "source_note", "Source: workbook reference",
              0.6, 6.9, 12, 0.4, size=10, italic=True, color=GRAY_TEXT)


# Dispatch by template id, mirroring the TS original.
BUILDERS = {
    "quarterly-review": build_quarterly_review,
}


def main() -> None:
    parser = argparse.ArgumentParser(description="Build branded .pptx templates from templates.json")
    parser.add_argument("--out", type=Path, default=_DEFAULT_OUT, help="output directory")
    args = parser.parse_args()
    args.out.mkdir(parents=True, exist_ok=True)

    defs = json.loads(_TEMPLATES_JSON.read_text(encoding="utf-8"))
    for d in defs:
        build = BUILDERS.get(d["id"])
        if build is None:
            raise SystemExit(f'no builder for template "{d["id"]}"')
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        build(prs)
        out_path = args.out / d["file"]
        prs.save(str(out_path))
        print(f"built {out_path} ({len(d['slides'])} slides)")


if __name__ == "__main__":
    main()
